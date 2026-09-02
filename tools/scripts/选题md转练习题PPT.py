# -*- coding: utf-8 -*-
r"""
从 选题.md 直接生成练习题 PPT（每题一页，贴题目图）。

适用于题目写在 选题.md 里、尚未入库（无题号）的情形：
解析 选题.md 的每个题块（题干 + A/B/C/D 选项 + 题目来源），
复用 练习题PPT.py 的 build_from_questions 渲图建页。

用法：
    python 选题md转练习题PPT.py <选题.md路径> <输出.pptx路径>
"""
import importlib.util, os, re, sys

sys.stdout.reconfigure(encoding="utf-8")
_HERE = os.path.dirname(os.path.abspath(__file__))


def _load(name, fname):
    spec = importlib.util.spec_from_file_location(name, os.path.join(_HERE, fname))
    m = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(m)
    return m


def parse_questions(md_text, img_dir=None):
    """按 ### 题块切分，提取 题干/选项/来源/图片，返回 dict 列表。
    img_dir: 图片所在目录的绝对路径。
    图片信息单独存在 figures 字段，不嵌入 LaTeX 内容（方便 PPT 中独立调整）。"""
    blocks = re.split(r"\n#{2,3}\s", md_text)
    out = []
    for b in blocks:
        if "【题干】" not in b:
            continue
        # 题干：从【题干】到第一个选项行 A. 之前
        stem_m = re.search(r"【题干】(.*?)(?=\n\s*A[\.、])", b, re.S)
        if not stem_m:
            # 没有选项的填空/解答题：题干取到下一个【标记】前
            stem_m = re.search(r"【题干】(.*?)(?=\n\s*【)", b, re.S)
            stem = stem_m.group(1).strip() if stem_m else ""
            opts = []
        else:
            stem = stem_m.group(1).strip()
            opts = []
            for letter in "ABCD":
                om = re.search(
                    rf"(?:^|\s){letter}[\.、]\s*(.*?)(?=\s[ABCD][\.、]|\n*【|$)", b, re.S
                )
                if om:
                    opts.append(om.group(1).strip())
        src_m = re.search(r"【题目来源】[ \t]*(.*)", b)
        src = src_m.group(1).strip() if src_m else ""

        # 提取图片引用（Obsidian ![[...]] + 标准 markdown ![alt](path)），单独存储
        figures = []

        def _add_figure(ref):
            """解析一个图片引用，加入 figures 列表。"""
            parts = ref.split("|")
            fname_raw = parts[0].strip()
            width_px = None
            if len(parts) > 1:
                try:
                    width_px = int(parts[1].strip())
                except ValueError:
                    pass
            # 绝对路径直接用，相对路径基于 img_dir 解析
            if os.path.isabs(fname_raw):
                img_path = fname_raw
            elif img_dir:
                img_path = os.path.join(img_dir, fname_raw)
            else:
                img_path = fname_raw
            if os.path.exists(img_path):
                figures.append({"path": img_path, "width_px": width_px})

        if img_dir:
            # Obsidian ![[...]]
            for m in re.finditer(r"!\[\[([^\]]+)\]\]", stem):
                _add_figure(m.group(1))
            stem = re.sub(r"!\[\[[^\]]+\]\]", "", stem)
            # 标准 markdown ![alt](path)
            for m in re.finditer(r"!\[[^\]]*\]\(([^)]+)\)", stem):
                _add_figure(m.group(1))
            stem = re.sub(r"!\[[^\]]*\]\([^)]+\)", "", stem)
        else:
            stem = re.sub(r"!\[\[[^\]]+\]\]", "", stem)
            stem = re.sub(r"!\[[^\]]*\]\([^)]+\)", "", stem)

        stem = re.sub(r"\n{3,}", "\n\n", stem).strip()  # 清理移除图片后多余的空行
        out.append({"content": stem, "options": opts, "source": src, "figures": figures})
    return out


def to_latex(q):
    src = q.get("source") or ""
    content = (q.get("content") or "").strip()
    text = f"（{src}）{content}" if src else content
    opts = q.get("options")
    if opts and len(opts) == 4:
        text += "\n\n\\choices" + "".join("{" + o + "}" for o in opts)
    return text


def main():
    md_path, out_path = sys.argv[1], sys.argv[2]
    pptx_mod = _load("lianxi", "练习题PPT.py")
    f = _load("biaoti_filler", "标题灌模板.py")
    from pptx import Presentation
    from pptx.util import Cm

    md = open(md_path, encoding="utf-8").read()
    img_dir = os.path.dirname(os.path.abspath(md_path))  # 图片在 md 同目录
    qs = parse_questions(md, img_dir=img_dir)
    if not qs:
        sys.exit(f"❌ 没解析出题目：{md_path}")

    prs = Presentation(pptx_mod.DEFAULT_TEMPLATE)
    template_slides = list(prs.slides)
    body = next((s for s in template_slides if f.T_BODY in f.slide_text(s)), None)
    if body is None:
        sys.exit("❌ 模板缺内容页")
    slides = []
    for _ in qs:
        s = f.duplicate_slide(prs, body)
        f.replace_tokens(s, {f.T_BODY: "题目讲解"})
        slides.append(s)
    imgs = f.render_question_images([to_latex(q) for q in qs])

    n_img = 0
    n_fig = 0
    for s, img, q in zip(slides, imgs, qs):
        # 贴文字渲染图
        if img is not None:
            f.insert_image(s, img)
            n_img += 1
        # 贴几何图形（单独放在文字图下方）
        for fig in q.get("figures", []):
            fig_path = fig["path"]
            if not os.path.exists(fig_path):
                continue
            # 计算位置：文字图底部 + 间距
            if img is not None:
                text_h_cm = img.height * Cm(2.5) / img.height if False else None
                # 用 EMU 算：insert_image 放在 Cm(2.5)，宽度 Cm(DISPLAY_WIDTH_CM)
                # 实际高度 = 原图高/宽 × 显示宽
                display_w = Cm(33.87)  # DISPLAY_WIDTH_CM
                ratio = img.height / img.width
                text_h = int(display_w * ratio)
                fig_top = Cm(2.5) + text_h + Cm(0.5)
            else:
                fig_top = Cm(2.5)
            # 图形宽度：Obsidian 标注的 px 换算（缩小一档，方便在 PPT 里调）
            if fig.get("width_px"):
                fig_w = Cm(max(3, min(10, fig["width_px"] / 70)))
            else:
                fig_w = Cm(5)
            s.shapes.add_picture(fig_path, Cm(1.5), fig_top, width=fig_w)
            n_fig += 1

    for s in template_slides:
        f.delete_slide(prs, s)
    os.makedirs(os.path.dirname(os.path.abspath(out_path)), exist_ok=True)
    prs.save(out_path)
    print(f"✅ {len(qs)} 题 / 渲图 {n_img} 张 / 贴图 {n_fig} 张 → {out_path}")
    if n_img < len(qs):
        print("  ⚠ 部分题目图未渲染（多半本地未装 MiKTeX）", file=sys.stderr)


if __name__ == "__main__":
    main()
