# -*- coding: utf-8 -*-
r"""
标题灌模板：把一份"标题 JSON"灌进微课 PPT 模板，生成 PPT 骨架。
- 普通页：只填标题（正文录课时手填）
- 例题页：可选地把题干用 LaTeX 渲成透明 PNG，自动贴到固定位置（依赖本地 MiKTeX）

模板只有两种页：
    封面   ：含记号 {{主标题}}、{{上标}}
    内容页 ：含记号 {{标题}}            ← 章节/知识点/例题/小结全用它

用法：
    python 标题灌模板.py --模板 微课PPT模板.pptx --数据 标题.json --输出 成品.pptx

JSON 结构：
    {
      "主标题": "复合函数的求导规则",
      "上标":   "导数及其应用",
      "页":     ["知识还原", "什么是复合函数？", "操作三步法"],
      "例题": [
        {"标题": "典型例题 1", "题目": "（教材例6）求 $y=(3x+5)^3$ 的导数。"},
        {"标题": "典型例题 2"}                          // 无"题目"=纯标题页
      ],
      "例题数": 0,            // 兼容写法：给整数则生成 N 张纯标题「典型例题 i」
      "小结":   true          // 默认 true，自动加「本讲小结」
    }
生成顺序：封面 → 页[] → 例题[] → 本讲小结
题目 LaTeX 格式见 tools/题目png生成工具/README.md（来源用全角括号、$公式$、\choices{}{}{}{}）。
"""
import argparse, copy, json, os, sys, tempfile
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Cm
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

T_COVER = "{{主标题}}"
T_BODY = "{{标题}}"

# 例题题目渲染/插入参数（统一在此调）
SLIDE_WIDTH_CM = 33.87                                          # 16:9 幻灯片宽（12192000 EMU）
IMG_SIDE_MARGIN_CM = 0                                          # 插入左右留白（0=整图铺满整页宽）
DISPLAY_WIDTH_CM = SLIDE_WIDTH_CM - 2 * IMG_SIDE_MARGIN_CM      # 插入宽=整图宽，1:1 不拉伸 → 字号真实
RENDER_WIDTH_CM = 28.22   # 传给渲染工具的"内容宽"；配合工具 MARGIN_LEFT/RIGHT 使整图≈PPT宽(33.87)、内容落在横线内、左对齐
FONT_PT = 22              # 题目字号（真实字号，不被拉伸改变）
DPI = 150                 # 题目图分辨率（投影足够；整图33.87cm@150dpi≈2000px）
IMG_TOP_CM = 2.5          # 题目图距顶

_render = None  # 延迟加载 compile_and_render


def _get_renderer():
    """返回批量渲染函数 compile_and_render_many（多题一次编译）。"""
    global _render
    if _render is None:
        tool_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "题目png生成工具")
        sys.path.insert(0, os.path.abspath(tool_dir))
        from latex_snippet_tool import compile_and_render_many
        _render = compile_and_render_many
    return _render


def _roman(n):
    """整数转罗马数字字符（典型例题编号用），如 1→Ⅰ、2→Ⅱ。"""
    nums = "ⅠⅡⅢⅣⅤⅥⅦⅧⅨⅩⅪⅫ"  # U+2160..U+216B
    return nums[n - 1] if 1 <= n <= len(nums) else str(n)


def slide_text(slide):
    return "\n".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)


def duplicate_slide(prs, source):
    new_slide = prs.slides.add_slide(source.slide_layout)
    for shp in list(new_slide.shapes):
        shp._element.getparent().remove(shp._element)
    for shp in source.shapes:
        new_slide.shapes._spTree.append(copy.deepcopy(shp._element))
    for rId, rel in source.part.rels.items():
        if rel.is_external or rel.reltype == RT.SLIDE_LAYOUT:
            continue
        if rel.reltype in (RT.IMAGE, RT.MEDIA):
            new_rId = new_slide.part.relate_to(rel._target, rel.reltype)
            for el in new_slide.shapes._spTree.iter():
                for attr in ("r:embed", "r:link"):
                    if el.get(qn(attr)) == rId:
                        el.set(qn(attr), new_rId)
    return new_slide


def replace_tokens(slide, mapping):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            full = "".join(r.text for r in para.runs)
            if "{{" not in full:
                continue
            new = full
            for token, val in mapping.items():
                new = new.replace(token, val)
            if new != full and para.runs:
                para.runs[0].text = new
                for r in para.runs[1:]:
                    r._r.getparent().remove(r._r)


def render_question_images(latex_list):
    """把多段题干 LaTeX **一次编译**渲成图片列表（顺序对应）。失败则返回全 None。"""
    if not latex_list:
        return []
    try:
        render_many = _get_renderer()
        return render_many(latex_list, width_cm=RENDER_WIDTH_CM, font_pt=FONT_PT, dpi=DPI)
    except Exception as e:
        print(f"  ⚠ 例题渲染失败，留空：{e}", file=sys.stderr)
        return [None] * len(latex_list)


def insert_image(slide, img):
    """把已渲染的 PIL 图片插到题目区固定位置。"""
    fd, png = tempfile.mkstemp(suffix=".png")
    os.close(fd)
    img.save(png)
    slide.shapes.add_picture(png, Cm(IMG_SIDE_MARGIN_CM), Cm(IMG_TOP_CM), width=Cm(DISPLAY_WIDTH_CM))
    os.remove(png)


def delete_slide(prs, slide):
    part = slide.part
    rId = next((rid for rid, rel in prs.part.rels.items()
                if not rel.is_external and rel._target is part), None)
    if rId is None:
        return
    for sldId in list(prs.slides._sldIdLst):
        if sldId.get(qn("r:id")) == rId:
            prs.slides._sldIdLst.remove(sldId)
            break
    prs.part.drop_rel(rId)


def build(template_path, data, out_path):
    prs = Presentation(template_path)
    template_slides = list(prs.slides)

    cover = body = None
    for s in template_slides:
        txt = slide_text(s)
        if T_COVER in txt:
            cover = s
        elif T_BODY in txt:
            body = s
    if cover is None or body is None:
        sys.exit(f"❌ 模板需含封面页({T_COVER})和内容页({T_BODY})各一张。")

    # 内容页序列：(标题, 题目LaTeX或None)
    pages = [(t, None) for t in data.get("页", [])]
    例题 = data.get("例题")
    if 例题:
        for e in 例题:
            pages.append((e.get("标题", "典型例题"), e.get("题目")))
    else:  # 兼容 例题数
        for i in range(1, int(data.get("例题数", 0)) + 1):
            pages.append((f"典型例题 {_roman(i)}", None))
    if data.get("小结", True):
        pages.append(("本讲小结", None))

    # 封面
    new_cover = duplicate_slide(prs, cover)
    replace_tokens(new_cover, {T_COVER: data.get("主标题", ""), "{{上标}}": data.get("上标", "")})
    # 内容页（先建页、收集待渲染例题 → 一次批量编译 → 贴图）
    n_img = 0
    pending = []  # [(slide, latex)]
    for title, latex in pages:
        s = duplicate_slide(prs, body)
        replace_tokens(s, {T_BODY: title})
        if latex:
            pending.append((s, latex))
    if pending:
        imgs = render_question_images([l for _, l in pending])
        for (s, _), img in zip(pending, imgs):
            if img is not None:
                insert_image(s, img)
                n_img += 1

    for s in template_slides:
        delete_slide(prs, s)

    os.makedirs(os.path.dirname(os.path.abspath(out_path)), exist_ok=True)
    prs.save(out_path)
    return len(prs.slides._sldIdLst), n_img


def main():
    ap = argparse.ArgumentParser(description="把标题 JSON 灌进微课 PPT 模板（例题可自动渲图）")
    ap.add_argument("--模板", dest="tpl", required=True)
    ap.add_argument("--数据", dest="data", required=True)
    ap.add_argument("--输出", dest="out", required=True)
    a = ap.parse_args()
    with open(a.data, encoding="utf-8") as f:
        data = json.load(f)
    n, n_img = build(a.tpl, data, a.out)
    print(f"✅ 已生成 {n} 页（其中 {n_img} 张例题图）→ {a.out}")


if __name__ == "__main__":
    main()
