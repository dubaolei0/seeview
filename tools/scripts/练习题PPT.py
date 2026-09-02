# -*- coding: utf-8 -*-
r"""
练习题 PPT：给一节课的练习题单独生成一份 PPT。

- 每道题一张内容页，小标题统一「题目讲解」，下面贴该题的题目图。
- 题目图 = 题库里的题（题干 + 选项 + 规范来源）用 LaTeX 渲成透明 PNG。
- 只用内容页，不出封面、不出小结。

复用：
- tools/scripts/标题灌模板.py 的复制页 / 填标题 / 渲图函数（不重复造）。
- tools/讲义生成工具/微课PPT模板.pptx 的内容页（含记号 {{标题}}）。
- 题源取自 master_database.jsonl + contributions.jsonl（与 MCP 同库）。

用法：
    python 练习题PPT.py --题号 u_6ae337e66b q_eee115290b ... --输出 导数练习题.pptx
    可选：--模板 <pptx>  --标题 题目讲解

依赖：python-pptx；渲图依赖本地 MiKTeX（无则该页留空标题、不报错）。
"""
import argparse, importlib.util, json, os, sys

sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation

_HERE = os.path.dirname(os.path.abspath(__file__))
_BASE = os.path.normpath(os.path.join(_HERE, "..", ".."))
_DB_DIR = os.path.join(_BASE, "knowledge", "高考题目", "题库")
MASTER_PATH = os.path.join(_DB_DIR, "master_database.jsonl")
TEXTBOOK_PATH = os.path.join(_DB_DIR, "textbook_questions.jsonl")
CONTRIB_PATH = os.path.join(_DB_DIR, "contributions.jsonl")
PATCH_PATH = os.path.join(_DB_DIR, "patches.jsonl")
DEFAULT_TEMPLATE = os.path.join(_BASE, "tools", "讲义生成工具", "微课PPT模板.pptx")


def load_db():
    """读题库为 id→题 字典。口径与 MCP 服务一致：
    master(高考拆题) → textbook(教材题) → contributions(用户新增)，同 id 后者覆盖前者；
    再叠加 patches.jsonl 的字段补丁（主库文件不变，可回溯）。"""
    idx = {}
    for path in (MASTER_PATH, TEXTBOOK_PATH, CONTRIB_PATH):
        if not os.path.exists(path):
            continue
        with open(path, encoding="utf-8") as f:
            for line in f:
                line = line.strip()
                if line:
                    q = json.loads(line)
                    idx[q["id"]] = q
    # 字段补丁：覆盖到已载入记录上
    if os.path.exists(PATCH_PATH):
        with open(PATCH_PATH, encoding="utf-8") as f:
            for line in f:
                line = line.strip()
                if not line:
                    continue
                patch = json.loads(line)
                target = idx.get(patch.get("target_id"))
                if target:
                    for k, v in (patch.get("fields") or {}).items():
                        target[k] = v
    return idx


def source_str(q):
    """与 MCP 的 _source_str 同口径：高考题/教材题 + ·拆解/·改编/·节选。"""
    t = q.get("type")
    if t == "textbook":
        base = q.get("source_exam") or "教材题"
    else:
        base = f"{q.get('source_year')}年{q.get('source_exam')}第{q.get('source_question_no')}题"
    variant = (q.get("source_variant") or "").strip().lstrip("·")
    if t in ("molecular", "atomic"):
        suffix = "·拆解"
    elif variant:
        suffix = f"·{variant}"
    else:
        suffix = ""
    return base + suffix


def to_latex(q):
    """把题转成 题目png生成工具 的输入格式：（来源）题干 + \\choices{}{}{}{}。"""
    src = source_str(q)
    content = (q.get("content") or "").replace("\\n", "\n").strip()
    text = f"（{src}）{content}"
    opts = q.get("options")
    if opts and q.get("question_format") in ("single_choice", "multiple_choice"):
        four = [str(o) for o in opts[:4]]
        text += "\n\n\\choices" + "".join("{" + o + "}" for o in four)
    return text


def _load_filler():
    """按路径加载中文文件名的 标题灌模板.py，复用其页操作函数。"""
    path = os.path.join(_HERE, "标题灌模板.py")
    spec = importlib.util.spec_from_file_location("biaoti_filler", path)
    m = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(m)
    return m


def build_from_questions(template_path, questions, out_path, title="题目讲解"):
    """核心：给一组题（dict 列表）生成练习题 PPT。供 CLI 与 MCP 工具共用。
    questions 顺序即页序；返回 (总页数, 成功渲染的题目图数)。"""
    f = _load_filler()
    prs = Presentation(template_path)
    template_slides = list(prs.slides)
    body = next((s for s in template_slides if f.T_BODY in f.slide_text(s)), None)
    if body is None:
        raise ValueError(f"模板缺内容页（需含记号 {f.T_BODY}）。")

    # 先建页 → 收集所有题目 LaTeX → 一次批量编译 → 贴图
    slides = []
    for q in questions:
        s = f.duplicate_slide(prs, body)
        f.replace_tokens(s, {f.T_BODY: title})
        slides.append(s)

    n_img = 0
    imgs = f.render_question_images([to_latex(q) for q in questions])
    for s, img in zip(slides, imgs):
        if img is not None:
            f.insert_image(s, img)
            n_img += 1

    for s in template_slides:
        f.delete_slide(prs, s)

    os.makedirs(os.path.dirname(os.path.abspath(out_path)), exist_ok=True)
    prs.save(out_path)
    return len(prs.slides._sldIdLst), n_img


def build(template_path, qids, out_path, title="题目讲解"):
    """CLI 入口：按题号从题库取题再生成。"""
    db = load_db()
    missing = [i for i in qids if i not in db]
    if missing:
        sys.exit(f"❌ 以下题号不存在于题库：{missing}")
    return build_from_questions(template_path, [db[i] for i in qids], out_path, title)


def main():
    ap = argparse.ArgumentParser(description="给一节课的练习题生成单独 PPT（每题一页，贴题目图）")
    ap.add_argument("--题号", dest="ids", nargs="+", required=True, help="题目 ID 列表（顺序即页序）")
    ap.add_argument("--输出", dest="out", required=True, help="输出 pptx 路径")
    ap.add_argument("--模板", dest="tpl", default=DEFAULT_TEMPLATE, help="PPT 模板（默认微课PPT模板.pptx）")
    ap.add_argument("--标题", dest="title", default="题目讲解", help="每页小标题（默认：题目讲解）")
    a = ap.parse_args()
    n, n_img = build(a.tpl, a.ids, a.out, a.title)
    print(f"✅ 已生成 {n} 页（{n_img} 张题目图）→ {a.out}")
    if n_img < len(a.ids):
        print("  ⚠ 部分题目图未渲染（多半是本地未装 MiKTeX）；标题页已就位，可装好后重跑。", file=sys.stderr)


if __name__ == "__main__":
    main()
